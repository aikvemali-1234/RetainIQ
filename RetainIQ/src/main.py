from fastapi import FastAPI, Request, Depends, Form, UploadFile, File, HTTPException
from fastapi.responses import RedirectResponse, HTMLResponse
from fastapi.staticfiles import StaticFiles
from fastapi.templating import Jinja2Templates
from starlette.middleware.sessions import SessionMiddleware

from sqlalchemy import create_engine, Column, Integer, String, Text, ForeignKey, DateTime
from sqlalchemy.orm import sessionmaker, declarative_base, relationship, Session as DBSession
from sqlalchemy.orm import Session

from passlib.context import CryptContext

import json
from sqlalchemy import func

import os
import shutil
from datetime import datetime
import random
import re
import json



from PyPDF2 import PdfReader
from docx import Document as DocxDocument
from pptx import Presentation

from openai import OpenAI  # OpenAI client

# ======================
# CONFIG & SETUP
# ======================

BASE_DIR = os.path.dirname(os.path.abspath(__file__))
UPLOAD_DIR = os.path.join(BASE_DIR, "uploads")
os.makedirs(UPLOAD_DIR, exist_ok=True)

app = FastAPI()

# Session middleware: used for login sessions
app.add_middleware(
    SessionMiddleware,
    secret_key="retainiq_123456789_super_secret_key_xyz",
    session_cookie="retainiq_session",
    same_site="lax",
    https_only=False,
)

app.mount("/static", StaticFiles(directory="static"), name="static")

templates = Jinja2Templates(directory="templates")

# Database
SQLALCHEMY_DATABASE_URL = "sqlite:///./retainiq.db"
engine = create_engine(SQLALCHEMY_DATABASE_URL, connect_args={"check_same_thread": False})
SessionLocal = sessionmaker(autocommit=False, autoflush=False, bind=engine)
Base = declarative_base()

# Password hashing
pwd_context = CryptContext(schemes=["pbkdf2_sha256"], deprecated="auto")


# 👇 HARD-CODED API KEY FOR HACKATHON DEMO (simple + reliable)
OPENAI_API_KEY = "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"

if not OPENAI_API_KEY:
    print("WARNING: OPENAI_API_KEY not set. AI quiz & flashcards will fall back.")
    client = None
else:
    client = OpenAI(api_key=OPENAI_API_KEY)
    print("✅ OPENAI_API_KEY loaded, AI features enabled.")




if OPENAI_API_KEY:
    client = OpenAI(api_key=OPENAI_API_KEY)
else:
    client = None
    print("WARNING: OPENAI_API_KEY not set. AI quiz generation will fall back to rule-based generator.")


# In-memory quiz state per user (for demo / hackathon)
QUIZ_STATE: dict[int, dict] = {}
FLASHCARDS_STATE: dict[tuple[int, int], dict] = {}  # key = (user_id, doc_id)


# ======================
# DB MODELS
# ======================

class User(Base):
    __tablename__ = "users"

    id = Column(Integer, primary_key=True, index=True)
    email = Column(String, unique=True, index=True)
    name = Column(String)
    password_hash = Column(String)
    total_points = Column(Integer, default=0)
    best_score = Column(Integer, default=0)
    badges = Column(Text, default="")  # comma-separated badge names
    linkedin_url = Column(String, default="")

    documents = relationship("Document", back_populates="user")
    attempts = relationship("QuizAttempt", back_populates="user")




class Document(Base):
    __tablename__ = "documents"

    id = Column(Integer, primary_key=True, index=True)
    user_id = Column(Integer, ForeignKey("users.id"))
    filename = Column(String)
    content = Column(Text)
    created_at = Column(DateTime, default=datetime.utcnow)   # <-- new

    user = relationship("User", back_populates="documents")
    attempts = relationship("QuizAttempt", back_populates="document")



class QuizAttempt(Base):
    __tablename__ = "quiz_attempts"

    id = Column(Integer, primary_key=True, index=True)
    user_id = Column(Integer, ForeignKey("users.id"))
    document_id = Column(Integer, ForeignKey("documents.id"))
    score = Column(Integer)
    total_questions = Column(Integer)
    created_at = Column(DateTime, default=datetime.utcnow)

    user = relationship("User", back_populates="attempts")
    document = relationship("Document", back_populates="attempts")


class QuizQuestionResult(Base):
    __tablename__ = "quiz_question_results"

    id = Column(Integer, primary_key=True, index=True)
    user_id = Column(Integer, ForeignKey("users.id"))
    document_id = Column(Integer, ForeignKey("documents.id"))
    attempt_id = Column(Integer, ForeignKey("quiz_attempts.id"))
    question_text = Column(Text)
    difficulty = Column(String)
    was_correct = Column(Integer)  # 1 = correct, 0 = wrong


class FlashcardProgress(Base):
    __tablename__ = "flashcard_progress"

    id = Column(Integer, primary_key=True, index=True)
    user_id = Column(Integer, ForeignKey("users.id"))
    document_id = Column(Integer, ForeignKey("documents.id"))
    front = Column(Text)   # flashcard front text
    back = Column(Text)    # flashcard back text
    known_count = Column(Integer, default=0)
    review_count = Column(Integer, default=0)
    last_seen = Column(DateTime, default=datetime.utcnow)


Base.metadata.create_all(bind=engine)

# ======================
# UTILS
# ======================

def get_db():
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()


def get_password_hash(password: str) -> str:
    return pwd_context.hash(password)


def verify_password(password: str, hashed: str) -> bool:
    return pwd_context.verify(password, hashed)


def get_current_user(request: Request, db: DBSession = Depends(get_db)):
    user_id = request.session.get("user_id")
    if not user_id:
        return None
    return db.query(User).filter(User.id == user_id).first()


def login_required(request: Request, db: DBSession = Depends(get_db)):
    user = get_current_user(request, db)
    if not user:
        return RedirectResponse(url="/login", status_code=302)
    return user


def extract_text_from_file(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    text = ""
    if ext == ".pdf":
        reader = PdfReader(file_path)
        for page in reader.pages:
            page_text = page.extract_text() or ""
            text += page_text + "\n"
    elif ext == ".docx":
        doc = DocxDocument(file_path)
        for para in doc.paragraphs:
            text += (para.text or "") + "\n"
    elif ext in [".ppt", ".pptx"]:
        prs = Presentation(file_path)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += (shape.text or "") + "\n"
    else:
        raise ValueError("Unsupported file type")

    text = re.sub(r"\n+", "\n", text).strip()

    if not text:
        raise ValueError(
            "Could not extract any text from this file. "
            "If it's a scanned PDF or image-based slide, try a text-based file."
        )
    return text


def simple_sentence_split(text: str):
    parts = re.split(r"[.!?]\s+", text)
    sentences = [p.strip() for p in parts if len(p.strip()) > 20]
    return sentences


def generate_quiz_from_text(text: str, difficulty: str, num_questions: int | None, auto_balance: bool = False):
    """
    Simple fallback quiz generator (no AI):
    - Split into sentences
    - Each question: "What is the main idea of this sentence?"
    """
    sentences = simple_sentence_split(text)

    if len(sentences) < 4:
        lines = [l.strip() for l in text.splitlines() if len(l.strip()) > 20]
        sentences = lines

    if not sentences:
        raise ValueError("Not enough meaningful text to generate quiz questions.")

    max_possible = len(sentences)
    if auto_balance:
        base = min(max_possible, 15)
        num_questions = max(5, base)
    if not num_questions:
        num_questions = min(10, max_possible)

    num_questions = min(num_questions, max_possible)

    all_sentences = sentences.copy()
    while len(all_sentences) < 4:
        all_sentences.extend(sentences)

    sampled = random.sample(sentences, num_questions)

    quiz_questions = []
    difficulties = []

    if auto_balance:
        third = max(1, num_questions // 3)
        difficulties = (["easy"] * third +
                        ["medium"] * third +
                        ["hard"] * (num_questions - 2 * third))
        random.shuffle(difficulties)
    else:
        difficulties = [difficulty] * num_questions

    for idx, sent in enumerate(sampled):
        correct = sent
        pool = [s for s in all_sentences if s != sent]
        if len(pool) < 3:
            pool = pool * 3
        distractors = random.sample(pool, 3)
        options = distractors + [correct]
        random.shuffle(options)
        correct_index = options.index(correct)

        quiz_questions.append({
            "id": idx,
            "question": f"What is the main idea of this sentence?\n\"{sent[:140]}...\"",
            "options": options,
            "correct_index": correct_index,
            "difficulty": difficulties[idx]
        })

    return quiz_questions


def generate_flashcards_from_text(text: str):
    """
    Simple flashcards:
    Front: "Explain: <first few words>"
    Back: full sentence
    """
    sentences = simple_sentence_split(text)
    if len(sentences) < 3:
        lines = [l.strip() for l in text.splitlines() if len(l.strip()) > 15]
        sentences = lines

    if not sentences:
        return []

    flashcards = []
    for idx, sent in enumerate(sentences[:50]):
        front = "Explain: " + " ".join(sent.split(" ")[:6]) + " ..."
        back = sent
        flashcards.append({
            "id": idx,
            "front": front,
            "back": back
        })
    return flashcards


def ai_generate_flashcards_from_text(text: str, max_cards: int = 30):
    ...
    if client is None:
        print("OPENAI client not initialized. Falling back to rule-based flashcards.")
        return generate_flashcards_from_text(text)

    # 🔹 SPEED: reduce text + #cards
    max_cards = min(max_cards, 20)
    material = text[:3000]


    material = text[:5000]

    prompt = f"""
You are an educational AI that creates flashcards to help students revise.

Given the following study material, generate up to {max_cards} flashcards.

Each flashcard should:
- Focus on one key concept or fact
- Have a short prompt on the front (question / term / fill-in-the-blank)
- Have a clear, correct answer or explanation on the back
- Be suitable for exam revision

Return ONLY valid JSON in this format:

[
  {{
    "front": "What is X?",
    "back": "X is ..."
  }},
  ...
]

Study material:
\"\"\"{material}\"\"\"
"""

    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "You are a helpful AI that generates study flashcards."},
                {"role": "user", "content": prompt}
            ]
        )

        raw_content = response.choices[0].message.content

        if "```" in raw_content:
            raw_content = raw_content.split("```")[1]
            raw_content = raw_content.replace("json", "", 1).strip()

        items = json.loads(raw_content)

        flashcards = []
        for idx, card in enumerate(items):
            front = card.get("front", "").strip()
            back = card.get("back", "").strip()
            if not front or not back:
                continue
            flashcards.append({
                "id": idx,
                "front": front,
                "back": back
            })

        if not flashcards:
            print("AI flashcards invalid/empty. Falling back to rule-based flashcards.")
            return generate_flashcards_from_text(text)

        return flashcards

    except Exception as e:
        print("AI Flashcard Generation Error:", e)
        return generate_flashcards_from_text(text)



def ai_generate_quiz_questions(text: str, difficulty: str, num_questions: int):
    """
    Use OpenAI (gpt-4o-mini) to generate MCQ questions from the given text.
    """
    if client is None:
        print("OPENAI client not initialized. Falling back to rule-based quiz generator.")
        return generate_quiz_from_text(text, difficulty or "mixed", num_questions, auto_balance=True)

    # 🔹 SPEED: use less text + fewer questions
    num_questions = min(num_questions, 8)  # hard limit to 8 questions per quiz
    material = text[:3000]  # send only first 3000 characters instead of 5000



    material = text[:5000]

    prompt = f"""
You are an educational AI that creates high-quality multiple-choice questions.

Read the following study material and generate {num_questions} questions.

Each question must:
- Be concept-based (NOT just random facts)
- Have exactly 4 options
- Have exactly 1 correct option
- Include a difficulty label: "easy", "medium", or "hard"

Difficulty setting requested: {difficulty}.
If difficulty is "mixed", create a balanced mix of easy, medium, and hard.

IMPORTANT:
Return ONLY valid JSON, with NO extra text, in this exact format:

[
  {{
    "question": "Question text here",
    "options": ["option A", "option B", "option C", "option D"],
    "correct_index": 0,
    "difficulty": "easy",
    "explanation": "Short explanation of the correct answer"
  }},
  ...
]

Study material:
\"\"\"{material}\"\"\"
"""

    try:
        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "You are a helpful AI that generates exam-style quizzes."},
                {"role": "user", "content": prompt}
            ]
        )

        raw_content = response.choices[0].message.content

        if "```" in raw_content:
            raw_content = raw_content.split("```")[1]
            raw_content = raw_content.replace("json", "", 1).strip()

        quiz_items = json.loads(raw_content)

        normalized = []
        for i, q in enumerate(quiz_items):
            normalized.append({
                "id": i,
                "question": q.get("question", "Untitled question"),
                "options": q.get("options", []),
                "correct_index": int(q.get("correct_index", 0)),
                "difficulty": q.get("difficulty", difficulty or "medium"),
                "explanation": q.get("explanation", "")
            })

        filtered = [q for q in normalized if len(q["options"]) == 4]
        if not filtered:
            print("AI returned invalid options; falling back to rule-based quiz generator.")
            return generate_quiz_from_text(text, difficulty or "mixed", num_questions, auto_balance=True)

        return filtered

    except Exception as e:
        print("AI Quiz Generation Error:", e)
        return generate_quiz_from_text(text, difficulty or "mixed", num_questions, auto_balance=True)


def award_badges(user: User, last_score: int):
    badges = set([b for b in user.badges.split(",") if b.strip()])
    if last_score > user.best_score:
        badges.add("New Personal Best")
    if user.total_points >= 100:
        badges.add("100+ Points Club")
    user.badges = ",".join(sorted(badges))


def compute_user_stats(db: DBSession, user_id: int) -> dict:
    """Compute core stats, badges and leaderboard position for a user."""
    # Quiz attempts & question results
    attempts = db.query(QuizAttempt).filter(
        QuizAttempt.user_id == user_id
    ).order_by(QuizAttempt.created_at.desc()).all()

    q_results = db.query(QuizQuestionResult).filter(
        QuizQuestionResult.user_id == user_id
    ).all()

    total_quizzes = len(attempts)
    total_questions = sum(a.total_questions for a in attempts) if attempts else 0
    total_score = sum(a.score for a in attempts) if attempts else 0
    avg_score = (total_score / total_quizzes) if total_quizzes else 0

    total_correct = sum(r.was_correct for r in q_results)
    accuracy = (total_correct / total_questions * 100) if total_questions else 0

    # Simple achievement badges
    badges: list[str] = []

    if total_quizzes >= 1:
        badges.append("First Quiz Completed")
    if total_quizzes >= 5:
        badges.append("Quiz Explorer (5+ quizzes)")
    if total_questions >= 50:
        badges.append("Question Warrior (50+ questions)")
    if accuracy >= 60:
        badges.append("Accuracy 60%+")
    if accuracy >= 80:
        badges.append("Sharp Mind 80%+")

    badge_count = len(badges)

    # Leaderboard position based on total_score (higher is better)
    score_subq = (
        db.query(
            QuizAttempt.user_id,
            func.sum(QuizAttempt.score).label("score_sum")
        )
        .group_by(QuizAttempt.user_id)
        .subquery()
    )

    leaderboard_position = None
    rows = db.query(score_subq.c.user_id, score_subq.c.score_sum).order_by(
        score_subq.c.score_sum.desc()
    ).all()

    for idx, row in enumerate(rows, start=1):
        if row.user_id == user_id:
            leaderboard_position = idx
            break

    return {
        "total_quizzes": total_quizzes,
        "total_questions": total_questions,
        "total_score": total_score,
        "avg_score": round(avg_score, 1),
        "accuracy": round(accuracy, 1),
        "badges": badges,
        "badge_count": badge_count,
        "leaderboard_position": leaderboard_position,
    }

def ai_generate_dashboard_insights(
    stats: dict,
    diff_stats: dict,
    hard_questions: list,
    review_cards: list,
) -> str:
    # TEMP: no external API calls, just a simple message
    if stats.get("total_quizzes", 0) == 0:
        return "Complete at least one quiz to unlock personalized study insights."

    return (
        f"AI insights are temporarily disabled. "
        f"Your overall accuracy is {stats.get('accuracy', 0)}%. "
        f"Try focusing on the topics where you got the most questions wrong."
    )

    # If there is almost no data, don't waste an API call
    if stats.get("total_quizzes", 0) == 0:
        return "Complete at least one quiz to unlock personalized AI study insights."

    try:
        payload = {
            "stats": stats,
            "by_difficulty": diff_stats,
            "hard_questions_count": len(hard_questions),
            "review_cards_count": len(review_cards),
        }

        prompt = f"""
You are a friendly study coach.

Here is a JSON summary of a learner's performance data:

{json.dumps(payload, default=str)}

Write a SHORT analysis (max 6 bullet points total):

- 2 bullet points about their strengths.
- 2 bullet points about their weaknesses or risk areas.
- 2 very specific suggestions of what they should revise next.

Use simple language, motivating tone. Do NOT repeat the numbers exactly,
just interpret them. Reply in plain text with bullet points.
"""

        response = client.chat.completions.create(
            model="gpt-4o-mini",
            messages=[
                {"role": "system", "content": "You are a concise, motivational learning coach."},
                {"role": "user", "content": prompt},
            ],
        )
        return response.choices[0].message.content.strip()

    except Exception as e:
        print("AI Dashboard Insights Error:", e)
        return (
            f"AI insights could not be generated right now. "
            f"For now, focus on improving your accuracy above {stats.get('accuracy', 0)}%."
        )



# ======================
# ROUTES
# ======================

@app.get("/", response_class=HTMLResponse)
async def root(request: Request):
    if request.session.get("user_id"):
        return RedirectResponse(url="/upload", status_code=302)
    return RedirectResponse(url="/login", status_code=302)


# ---- AUTH ----

@app.get("/signup", response_class=HTMLResponse)
async def signup_page(request: Request):
    return templates.TemplateResponse("signup.html", {"request": request})


@app.post("/signup")
async def signup(
    request: Request,
    name: str = Form(...),
    email: str = Form(...),
    password: str = Form(...),
    db: DBSession = Depends(get_db)
):
    existing = db.query(User).filter(User.email == email).first()
    if existing:
        return templates.TemplateResponse(
            "signup.html",
            {"request": request, "error": "Email already registered"}
        )
    user = User(
        name=name,
        email=email,
        password_hash=get_password_hash(password)
    )
    db.add(user)
    db.commit()
    db.refresh(user)
    request.session["user_id"] = user.id
    return RedirectResponse(url="/dashboard", status_code=302)


@app.get("/login", response_class=HTMLResponse)
async def login_page(request: Request):
    return templates.TemplateResponse("login.html", {"request": request})


@app.post("/login")
async def login(
    request: Request,
    email: str = Form(...),
    password: str = Form(...),
    db: DBSession = Depends(get_db)
):
    user = db.query(User).filter(User.email == email).first()
    if not user or not verify_password(password, user.password_hash):
        return templates.TemplateResponse(
            "login.html",
            {"request": request, "error": "Invalid email or password"}
        )
    request.session["user_id"] = user.id
    return RedirectResponse(url="/dashboard", status_code=302)


@app.get("/logout")
async def logout(request: Request):
    request.session.clear()
    return RedirectResponse(url="/login", status_code=302)


# ---- UPLOAD & MODE SELECT ----

@app.get("/upload", response_class=HTMLResponse)
async def upload_page(
    request: Request,
    user: User = Depends(login_required)
):
    return templates.TemplateResponse(
        "upload.html",
        {"request": request, "user": user}
    )



@app.post("/upload")
async def upload_file(
    request: Request,
    file: UploadFile = File(...),
    db: DBSession = Depends(get_db),
    user: User = Depends(login_required)
):
    ext = os.path.splitext(file.filename)[1].lower()
    if ext not in [".pdf", ".docx", ".ppt", ".pptx"]:
        return templates.TemplateResponse(
            "upload.html",
            {
                "request": request,
                "user": user,
                "error": "Only PDF, DOCX, PPT/PPTX files are supported."
            }
        )

    file_path = os.path.join(UPLOAD_DIR, file.filename)
    with open(file_path, "wb") as buffer:
        shutil.copyfileobj(file.file, buffer)

    try:
        text = extract_text_from_file(file_path)
    except Exception as e:
        return templates.TemplateResponse(
            "upload.html",
            {
                "request": request,
                "user": user,
                "error": f"Failed to read file: {e}"
            }
        )

    doc = Document(
        user_id=user.id,
        filename=file.filename,
        content=text
    )
    db.add(doc)
    db.commit()
    db.refresh(doc)

    return RedirectResponse(url=f"/mode/{doc.id}", status_code=302)


@app.get("/mode/{doc_id}", response_class=HTMLResponse)
async def mode_select(
    request: Request,
    doc_id: int,
    db: DBSession = Depends(get_db),
    user: User = Depends(login_required)
):
    doc = db.query(Document).filter(
        Document.id == doc_id,
        Document.user_id == user.id
    ).first()
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
    return templates.TemplateResponse(
        "mode_select.html",
        {"request": request, "user": user, "doc": doc}
    )


# ---- QUIZ FLOW ----

@app.get("/quiz/settings/{doc_id}", response_class=HTMLResponse)
async def quiz_settings(
    request: Request,
    doc_id: int,
    db: DBSession = Depends(get_db),
    user: User = Depends(login_required)
):
    doc = db.query(Document).filter(
        Document.id == doc_id,
        Document.user_id == user.id
    ).first()
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")
    return templates.TemplateResponse(
        "quiz_settings.html",
        {"request": request, "user": user, "doc": doc}
    )


@app.post("/quiz/start/{doc_id}")
async def quiz_start(
    request: Request,
    doc_id: int,
    difficulty: str = Form(...),
    question_count_mode: str = Form(...),
    db: DBSession = Depends(get_db),
    user: User = Depends(login_required)
):
    print("QUIZ START ROUTE HIT — DEBUG")
    print("Difficulty:", difficulty)
    print("Question mode:", question_count_mode)

    doc = db.query(Document).filter(
        Document.id == doc_id,
        Document.user_id == user.id
    ).first()
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")

    auto_balance = False
    num_questions = None

    if difficulty == "custom_doc":
        auto_balance = True
    else:
        if question_count_mode == "10-15":
            num_questions = random.randint(10, 15)
        elif question_count_mode == "15-20":
            num_questions = random.randint(15, 20)
        elif question_count_mode == "custom_full":
            num_questions = min(
                25,
                max(10, len(simple_sentence_split(doc.content)) // 2)
            )

    if not num_questions:
        if question_count_mode == "10-15":
            num_questions = random.randint(10, 15)
        elif question_count_mode == "15-20":
            num_questions = random.randint(15, 20)
        else:
            num_questions = 10

    effective_difficulty = difficulty if difficulty != "custom_doc" else "mixed"

    quiz_questions = ai_generate_quiz_questions(
        doc.content,
        effective_difficulty,
        num_questions
    )

    QUIZ_STATE[user.id] = {
        "doc_id": doc.id,
        "questions": quiz_questions,
        "current_index": 0,
        "score": 0,
        "answers": []
    }

    print("QUIZ STATE SET IN MEMORY — DEBUG")

    return RedirectResponse(url="/quiz/play", status_code=302)


@app.get("/quiz/play", response_class=HTMLResponse)
async def quiz_play(
    request: Request,
    user: User = Depends(login_required)
):
    print("QUIZ PLAY ROUTE HIT — DEBUG")

    quiz = QUIZ_STATE.get(user.id)
    print("QUIZ FROM MEMORY:", "FOUND" if quiz else "MISSING")

    if not quiz:
        return RedirectResponse(url="/upload", status_code=302)

    idx = quiz["current_index"]
    questions = quiz["questions"]

    if idx >= len(questions):
        return RedirectResponse(url="/quiz/result", status_code=302)

    question = questions[idx]

    return templates.TemplateResponse(
        "quiz.html",
        {
            "request": request,
            "user": user,
            "index": idx,
            "total": len(questions),
            "question": question,
        }
    )


@app.post("/quiz/answer")
async def quiz_answer(
    request: Request,
    selected_option: int = Form(...),
    user: User = Depends(login_required)
):
    quiz = QUIZ_STATE.get(user.id)
    if not quiz:
        return RedirectResponse(url="/upload", status_code=302)

    idx = quiz["current_index"]
    questions = quiz["questions"]
    question = questions[idx]

    correct = (selected_option == question["correct_index"])
    if correct:
        quiz["score"] += 10

    quiz["answers"].append({
        "question": question,
        "selected": selected_option,
        "correct": correct
    })
    quiz["current_index"] += 1

    QUIZ_STATE[user.id] = quiz

    if quiz["current_index"] >= len(questions):
        return RedirectResponse(url="/quiz/result", status_code=302)
    else:
        return RedirectResponse(url="/quiz/play", status_code=302)


@app.get("/quiz/result", response_class=HTMLResponse)
async def quiz_result(
    request: Request,
    db: DBSession = Depends(get_db),
    user: User = Depends(login_required)
):
    quiz = QUIZ_STATE.get(user.id)
    if not quiz:
        return RedirectResponse(url="/upload", status_code=302)

    total_questions = len(quiz["questions"])
    score = quiz["score"]
    answers = quiz["answers"]
    doc_id = quiz["doc_id"]

    attempt = QuizAttempt(
        user_id=user.id,
        document_id=doc_id,
        score=score,
        total_questions=total_questions
    )
    db.add(attempt)
    db.commit()
    db.refresh(attempt)

    for ans in answers:
        q = ans["question"]
        result = QuizQuestionResult(
            user_id=user.id,
            document_id=doc_id,
            attempt_id=attempt.id,
            question_text=q.get("question", ""),
            difficulty=q.get("difficulty", ""),
            was_correct=1 if ans["correct"] else 0,
        )
        db.add(result)

    user.total_points += score
    if score > user.best_score:
        user.best_score = score
    award_badges(user, score)

    db.commit()

    QUIZ_STATE.pop(user.id, None)

    return templates.TemplateResponse(
        "quiz_result.html",
        {
            "request": request,
            "user": user,
            "score": score,
            "total_questions": total_questions,
            "answers": answers
        }
    )


# ---- FLASHCARDS FLOW ----

@app.get("/flashcards/{doc_id}", response_class=HTMLResponse)
async def flashcards_view(
    request: Request,
    doc_id: int,
    db: DBSession = Depends(get_db),
    user: User = Depends(login_required),
):
    # Load doc and ownership check
    doc = db.query(Document).filter(Document.id == doc_id, Document.user_id == user.id).first()
    if not doc:
        return RedirectResponse("/dashboard", status_code=302)

    sess = request.session or {}
    # Ensure top-level container exists
    flash_state = sess.get("flashcards", {})

    # If no flashcards generated for this doc, generate and save in session
    key = str(doc_id)
    if key not in flash_state:
        # generate flashcards (AI or fallback)
        try:
            cards = ai_generate_flashcards_from_text(doc.content, max_cards=20)
        except Exception:
            cards = generate_flashcards_from_text(doc.content)
        # ensure list (never None)
        cards = cards or []
        flash_state[key] = {"cards": cards, "index": 0, "show_back": False}
        request.session["flashcards"] = flash_state  # persist

    # load current state for doc
    state = flash_state.get(key, {"cards": [], "index": 0, "show_back": False})
    cards = state.get("cards", [])
    idx = int(state.get("index", 0) or 0)
    idx = max(0, min(idx, max(0, len(cards) - 1)))

    card = cards[idx] if cards else None

    return templates.TemplateResponse(
        "flashcards.html",
        {
            "request": request,
            "user": user,
            "card": card,
            "index": idx + 1,
            "total": len(cards),
            "show_back": bool(state.get("show_back", False)),
            "doc_id": doc_id,
        },
    )


from fastapi import Form

@app.post("/flashcards/action/{doc_id}")
async def flashcards_action(
    request: Request,
    doc_id: int,
    action: str = Form(...),
    db: DBSession = Depends(get_db),
    user: User = Depends(login_required),
):
    # basic ownership check (safe guard)
    doc = db.query(Document).filter(Document.id == doc_id, Document.user_id == user.id).first()
    if not doc:
        return RedirectResponse("/dashboard", status_code=302)

    flash_state = request.session.get("flashcards", {})  # dict
    key = str(doc_id)
    if key not in flash_state:
        # nothing to act upon — regenerate and redirect to view
        return RedirectResponse(f"/flashcards/{doc_id}", status_code=302)

    state = flash_state[key]
    cards = state.get("cards", [])
    idx = int(state.get("index", 0) or 0)
    idx = max(0, min(idx, max(0, len(cards) - 1)))

    # actions
    if action == "flip":
        state["show_back"] = not bool(state.get("show_back", False))

    elif action == "next":
        state["index"] = min(idx + 1, max(0, len(cards) - 1))
        state["show_back"] = False

    elif action == "prev":
        state["index"] = max(idx - 1, 0)
        state["show_back"] = False

    # persist back to session (must overwrite top-level flashcards)
    flash_state[key] = state
    request.session["flashcards"] = flash_state

    # always redirect to GET so the page is rendered normally
    return RedirectResponse(f"/flashcards/{doc_id}", status_code=302)


# ---- PROFILE & DASHBOARD ----

@app.get("/profile", response_class=HTMLResponse)
async def profile(
    request: Request,
    db: DBSession = Depends(get_db),
    user: User = Depends(login_required)
):
    attempts = db.query(QuizAttempt).filter(
        QuizAttempt.user_id == user.id
    ).order_by(QuizAttempt.created_at.desc()).all()

    top_users = db.query(User).order_by(User.total_points.desc()).limit(5).all()

    badges = [b for b in user.badges.split(",") if b.strip()]

    return templates.TemplateResponse(
        "profile.html",
        {
            "request": request,
            "user": user,
            "attempts": attempts,
            "badges": badges,
            "top_users": top_users
        }
    )


@app.post("/profile/linkedin")
async def update_linkedin(
    request: Request,
    linkedin_url: str = Form(...),
    db: DBSession = Depends(get_db),
    user: User = Depends(login_required)
):
    user.linkedin_url = linkedin_url
    db.commit()
    return RedirectResponse(url="/profile", status_code=302)



@app.get("/dashboard")
def dashboard(request: Request, session: Session = Depends(get_db)):

    # --- Ensure user logged in ---
    user_id = request.session.get("user_id")
    if not user_id:
        return RedirectResponse("/login", status_code=302)

    user = session.get(User, user_id)

    # --- Safe user display name ---
    display_name = (
        getattr(user, "full_name", None)
        or getattr(user, "name", None)
        or getattr(user, "username", None)
        or user.email.split("@")[0]
    )

    # --- Load user's documents ---
    documents = (
        session.query(Document)
        .filter(Document.user_id == user_id)
        .order_by(Document.id.desc())
        .all()
    )

    # --- Compute simple stats ---
    total_quizzes = (
        session.query(QuizAttempt)
        .filter(QuizAttempt.user_id == user_id)
        .count()
    )

    total_flashcards = (
        session.query(FlashcardProgress)
        .filter(FlashcardProgress.user_id == user_id)
        .count()
    )

    return templates.TemplateResponse(
        "dashboard.html",
        {
            "request": request,
            "display_name": display_name,
            "documents": documents,
            "total_quizzes": total_quizzes,
            "total_flashcards": total_flashcards,
        },
    )


@app.get("/dashboard/doc/{doc_id}")
def dashboard_doc(doc_id: int, request: Request, session: Session = Depends(get_db)):
    user_id = request.session.get("user_id")
    if not user_id:
        return RedirectResponse("/login", status_code=302)

    doc = (
        session.query(Document)
        .filter(Document.id == doc_id, Document.user_id == user_id)
        .first()
    )

    if not doc:
        return RedirectResponse("/dashboard", status_code=302)

    # --- Quiz attempts for this document ---
    attempts = (
        session.query(QuizAttempt)
        .filter(QuizAttempt.document_id == doc_id, QuizAttempt.user_id == user_id)
        .order_by(QuizAttempt.created_at.desc())
        .all()
    )

    # --- Flashcard progress for this document ---
    flashcards = (
        session.query(FlashcardProgress)
        .filter(FlashcardProgress.document_id == doc_id, FlashcardProgress.user_id == user_id)
        .all()
    )

    # --- Basic stats ---
    total_questions = sum(a.total_questions or 0 for a in attempts)
    total_score = sum(a.score or 0 for a in attempts)

    accuracy = 0
    if total_questions > 0:
        accuracy = round((total_score / (total_questions * 10)) * 100, 1)

    return templates.TemplateResponse(
        "document_view.html",
        {
            "request": request,
            "doc": doc,
            "attempts": attempts,
            "flashcards": flashcards,
            "total_questions": total_questions,
            "total_score": total_score,
            "accuracy": accuracy,
        },
    )


from fastapi import Form

@app.post("/documents/{doc_id}/rename")
async def rename_document(
    doc_id: int,
    new_name: str = Form(...),
    db: DBSession = Depends(get_db),
    user: User = Depends(login_required),
):
    doc = (
        db.query(Document)
        .filter(Document.id == doc_id, Document.user_id == user.id)
        .first()
    )
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")

    if hasattr(doc, "title"):
        doc.title = new_name
    else:
        doc.filename = new_name

    db.commit()
    return RedirectResponse("/dashboard", status_code=302)


@app.post("/documents/{doc_id}/delete")
async def delete_document(
    doc_id: int,
    db: DBSession = Depends(get_db),
    user: User = Depends(login_required),
):
    doc = (
        db.query(Document)
        .filter(Document.id == doc_id, Document.user_id == user.id)
        .first()
    )
    if not doc:
        raise HTTPException(status_code=404, detail="Document not found")

    # delete quiz question results & attempts
    attempt_ids = [
        a.id
        for a in db.query(QuizAttempt)
        .filter(QuizAttempt.user_id == user.id, QuizAttempt.document_id == doc_id)
        .all()
    ]
    if attempt_ids:
        db.query(QuizQuestionResult).filter(
            QuizQuestionResult.attempt_id.in_(attempt_ids)
        ).delete(synchronize_session=False)

    db.query(QuizAttempt).filter(
        QuizAttempt.user_id == user.id, QuizAttempt.document_id == doc_id
    ).delete(synchronize_session=False)

    # delete flashcard progress
    db.query(FlashcardProgress).filter(
        FlashcardProgress.user_id == user.id,
        FlashcardProgress.document_id == doc_id,
    ).delete(synchronize_session=False)

    # finally delete the document
    db.delete(doc)
    db.commit()

    return RedirectResponse("/dashboard", status_code=302)
